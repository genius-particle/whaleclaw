"""Post-processing for generated PPTX files.

Fixes common issues that LLM-generated scripts leave behind:
  1. Large opaque rectangles covering images (Z-order fix)
  2. Images covering text boxes (Z-order fix)
  3. Shapes overflowing the slide canvas (boundary clamp)
  4. Image crops that cut off faces (face-aware re-crop)
"""

from __future__ import annotations

import logging
from pathlib import Path

log = logging.getLogger(__name__)

# Minimum area fraction to consider a rectangle "large"
_LARGE_SHAPE_AREA_RATIO = 0.25

_RATIO_DIFF_THRESHOLD = 0.08


def fix_pptx(path: str | Path) -> bool:
    """Apply all post-processing fixes to a .pptx file in-place.

    Returns True if the file was modified, False otherwise.
    """
    path = Path(path)
    if not path.exists() or path.suffix.lower() != ".pptx":
        return False

    try:
        from pptx import Presentation
        from pptx.util import Emu
    except ImportError:
        return False

    try:
        prs = Presentation(str(path))
    except Exception:
        log.warning("pptx_postprocess: failed to open %s", path)
        return False

    sw = prs.slide_width
    sh = prs.slide_height
    slide_area = sw * sh
    modified = False

    for slide in prs.slides:
        sp_tree = slide.shapes._spTree  # type: ignore[attr-defined]
        children = list(sp_tree)

        image_indices: list[int] = []
        large_rect_indices: list[int] = []

        for idx, child in enumerate(children):
            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            is_picture = tag == "pic"
            is_sp = tag == "sp"

            if is_picture:
                image_indices.append(idx)
            elif is_sp:
                xfrm = child.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm")
                if xfrm is not None:
                    off = xfrm.find("{http://schemas.openxmlformats.org/drawingml/2006/main}off")
                    ext = xfrm.find("{http://schemas.openxmlformats.org/drawingml/2006/main}ext")
                    if off is not None and ext is not None:
                        cx = int(ext.get("cx", "0"))
                        cy = int(ext.get("cy", "0"))
                        area = cx * cy
                        if area >= slide_area * _LARGE_SHAPE_AREA_RATIO:
                            large_rect_indices.append(idx)

        if image_indices and large_rect_indices:
            min_img_idx = min(image_indices)
            rects_above_images = [i for i in large_rect_indices if i > min_img_idx]

            for rect_idx in sorted(rects_above_images, reverse=True):
                element = children[rect_idx]
                solid = element.find(
                    ".//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill"
                )
                if solid is None:
                    continue
                alpha_elem = solid.find(
                    "{http://schemas.openxmlformats.org/drawingml/2006/main}alpha"
                )
                if alpha_elem is not None:
                    val = int(alpha_elem.get("val", "100000"))
                    if val < 60000:
                        continue

                sp_tree.remove(element)
                insert_pos = 0
                for i, ch in enumerate(sp_tree):
                    ch_tag = ch.tag.split("}")[-1] if "}" in ch.tag else ch.tag
                    if ch_tag in ("cNvGrpSpPr", "grpSpPr", "nvGrpSpPr"):
                        insert_pos = i + 1
                    else:
                        break
                sp_tree.insert(insert_pos, element)
                modified = True

        if _fix_pictures_over_text(slide):
            modified = True
        _clamp_shapes(slide, sw, sh)
        if _fix_face_crops(slide):
            modified = True

    if modified:
        try:
            prs.save(str(path))
        except Exception:
            log.warning("pptx_postprocess: failed to save %s", path)
            return False

    return modified


_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_OVERLAP_THRESHOLD = 0.25


def _get_xfrm_rect(
    element: object,
) -> tuple[int, int, int, int] | None:
    """Extract (left, top, right, bottom) in EMU from an sp/pic element."""
    xfrm = element.find(f".//{{{_A}}}xfrm")  # type: ignore[union-attr]
    if xfrm is None:
        return None
    off = xfrm.find(f"{{{_A}}}off")
    ext = xfrm.find(f"{{{_A}}}ext")
    if off is None or ext is None:
        return None
    x = int(off.get("x", "0"))
    y = int(off.get("y", "0"))
    cx = int(ext.get("cx", "0"))
    cy = int(ext.get("cy", "0"))
    if cx <= 0 or cy <= 0:
        return None
    return (x, y, x + cx, y + cy)


def _overlap_fraction(
    r1: tuple[int, int, int, int],
    r2: tuple[int, int, int, int],
) -> float:
    """Fraction of r2's area overlapped by r1."""
    ix0 = max(r1[0], r2[0])
    iy0 = max(r1[1], r2[1])
    ix1 = min(r1[2], r2[2])
    iy1 = min(r1[3], r2[3])
    if ix0 >= ix1 or iy0 >= iy1:
        return 0.0
    inter = (ix1 - ix0) * (iy1 - iy0)
    area2 = (r2[2] - r2[0]) * (r2[3] - r2[1])
    return inter / area2 if area2 > 0 else 0.0


def _has_text_content(element: object) -> bool:
    """True if element contains non-whitespace text."""
    for t_elem in element.findall(f".//{{{_A}}}t"):  # type: ignore[union-attr]
        if t_elem.text and t_elem.text.strip():
            return True
    return False


def _fix_pictures_over_text(slide: object) -> bool:
    """Move pictures below text boxes when they overlap.

    In LLM-generated PPTX, images are often inserted after text boxes,
    placing them above text in Z-order and visually covering the text.
    """
    sp_tree = slide.shapes._spTree  # type: ignore[attr-defined]
    children = list(sp_tree)
    changed = False

    pic_entries: list[tuple[int, object, tuple[int, int, int, int]]] = []
    text_entries: list[tuple[int, object, tuple[int, int, int, int]]] = []

    for idx, child in enumerate(children):
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        rect = _get_xfrm_rect(child)
        if rect is None:
            continue
        if tag == "pic":
            pic_entries.append((idx, child, rect))
        elif tag == "sp" and _has_text_content(child):
            text_entries.append((idx, child, rect))

    if not pic_entries or not text_entries:
        return False

    pics_to_move: list[object] = []
    for pic_idx, pic_el, pic_rect in pic_entries:
        for txt_idx, _txt_el, txt_rect in text_entries:
            if pic_idx <= txt_idx:
                continue
            frac = _overlap_fraction(pic_rect, txt_rect)
            if frac >= _OVERLAP_THRESHOLD:
                pics_to_move.append(pic_el)
                break

    if not pics_to_move:
        return False

    first_content_pos = 0
    for i, ch in enumerate(sp_tree):
        ch_tag = ch.tag.split("}")[-1] if "}" in ch.tag else ch.tag
        if ch_tag in ("cNvGrpSpPr", "grpSpPr", "nvGrpSpPr"):
            first_content_pos = i + 1
        else:
            break

    for pic_el in pics_to_move:
        sp_tree.remove(pic_el)
        sp_tree.insert(first_content_pos, pic_el)
        changed = True

    return changed


def _clamp_shapes(slide: object, sw: int, sh: int) -> None:
    """Ensure no shape overflows the slide canvas and fix text overflow."""
    try:
        for shape in slide.shapes:  # type: ignore[attr-defined]
            if shape.left < 0:
                shape.width = max(0, shape.width + shape.left)
                shape.left = 0
            if shape.top < 0:
                shape.height = max(0, shape.height + shape.top)
                shape.top = 0
            if shape.left + shape.width > sw:
                shape.width = sw - shape.left
            if shape.top + shape.height > sh:
                shape.height = sh - shape.top

            _fix_textbox_overflow(shape, sw)
    except Exception:
        pass


def _fix_textbox_overflow(shape: object, slide_width: int) -> None:
    """Force word-wrap on text frames to prevent text overflow.

    Also applies TEXT_TO_FIT_SHAPE auto-size so that if text still doesn't
    fit after wrapping, the font shrinks rather than overflowing.
    """
    try:
        from pptx.enum.text import MSO_AUTO_SIZE  # type: ignore[import-untyped]
    except ImportError:
        return

    tf = getattr(shape, "text_frame", None)
    if tf is None:
        return
    text = getattr(tf, "text", "")
    if not text or not text.strip():
        return

    try:
        tf.word_wrap = True
    except Exception:
        pass

    try:
        cur = getattr(tf, "auto_size", None)
        if cur is None or cur == MSO_AUTO_SIZE.NONE:
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass

    right_edge = getattr(shape, "left", 0) + getattr(shape, "width", 0)
    if right_edge > slide_width:
        try:
            shape.width = slide_width - shape.left  # type: ignore[attr-defined]
        except Exception:
            pass


_IMG_SUFFIXES = {".jpg", ".jpeg", ".png", ".webp", ".bmp", ".tiff"}
_ORIGIN_RECENCY_SEC = 1800


def _origin_search_dirs() -> list[Path]:
    dirs = [Path("/tmp")]
    try:
        from whaleclaw.config.paths import WHALECLAW_HOME
        dl = WHALECLAW_HOME / "downloads"
        if dl.is_dir():
            dirs.insert(0, dl)
    except Exception:
        home_dl = Path.home() / ".whaleclaw" / "downloads"
        if home_dl.is_dir():
            dirs.insert(0, home_dl)
    return dirs


def _build_origin_index() -> dict[str, Path]:
    """Build a filename→path index of recent image files in /tmp."""
    import time

    cutoff = time.time() - _ORIGIN_RECENCY_SEC
    index: dict[str, Path] = {}
    skip_patterns = ("__cropped_tmp__", "__resized__", "__thumb__")
    for d in _origin_search_dirs():
        try:
            for p in d.iterdir():
                if p.suffix.lower() not in _IMG_SUFFIXES:
                    continue
                name_lower = p.name.lower()
                if any(pat in name_lower for pat in skip_patterns):
                    continue
                try:
                    if p.stat().st_mtime < cutoff:
                        continue
                except Exception:
                    continue
                index[name_lower] = p
        except Exception:
            continue
    return index


def _color_hist(img: object) -> object:
    """Compute a normalised color histogram for matching."""
    import cv2

    hsv = cv2.cvtColor(img, cv2.COLOR_BGR2HSV)  # type: ignore[arg-type]
    hist = cv2.calcHist([hsv], [0, 1], None, [16, 16], [0, 180, 0, 256])  # type: ignore[arg-type,list-item]
    cv2.normalize(hist, hist)  # type: ignore[arg-type]
    return hist


def _find_original_image(
    blob_img: object,
    origin_index: dict[str, Path],
    img_filename: str,
) -> Path | None:
    """Try to find the original (larger) source image for a PPT blob.

    Strategy 1: match by filename stem.
    Strategy 2: match by color histogram correlation.
    """
    import cv2

    stem = Path(img_filename).stem.lower() if img_filename else ""
    if stem and stem != "image":
        for name, path in origin_index.items():
            if Path(name).stem.lower() == stem:
                return path
        for name, path in origin_index.items():
            if stem in Path(name).stem.lower() or Path(name).stem.lower() in stem:
                return path

    blob_h, blob_w = blob_img.shape[:2]  # type: ignore[union-attr]
    blob_pixels = blob_w * blob_h
    blob_hist = _color_hist(blob_img)
    best_path: Path | None = None
    best_score = -1.0

    for _name, path in origin_index.items():
        orig = cv2.imread(str(path), cv2.IMREAD_COLOR)
        if orig is None:
            continue
        oh, ow = orig.shape[:2]
        if ow * oh <= blob_pixels:
            continue
        orig_hist = _color_hist(orig)
        score = cv2.compareHist(blob_hist, orig_hist, cv2.HISTCMP_CORREL)  # type: ignore[arg-type]
        if score > best_score:
            best_score = score
            best_path = path

    if best_score >= 0.6:
        return best_path
    return None


def _fix_face_crops(slide: object) -> bool:
    """Replace distorted/badly-cropped images with face-aware cropped versions.

    Two strategies:
      1. If blob ratio != box ratio → crop the blob directly.
      2. If blob was pre-resized by PIL (ratio matches but image is small),
         try to find the original larger image in /tmp and crop from that.
    """
    changed = False
    try:
        import cv2
        import numpy as np
    except ImportError:
        return False

    from whaleclaw.utils.image_crop import detect_face_info, smart_crop_box

    origin_index = _build_origin_index()

    for shape in slide.shapes:  # type: ignore[attr-defined]
        if getattr(shape, "shape_type", None) != 13:
            continue

        try:
            blob = shape.image.blob
            img_filename = shape.image.filename or ""
        except Exception:
            continue

        box_w = shape.width
        box_h = shape.height
        if not box_w or not box_h:
            continue
        box_ratio = box_w / box_h

        arr = np.frombuffer(blob, dtype=np.uint8)
        img = cv2.imdecode(arr, cv2.IMREAD_COLOR)
        if img is None:
            continue
        ih, iw = img.shape[:2]
        img_ratio = iw / ih
        ratio_diff = abs(img_ratio - box_ratio)

        source_img = img
        source_iw, source_ih = iw, ih

        if ratio_diff < _RATIO_DIFF_THRESHOLD:
            origin_path = _find_original_image(img, origin_index, img_filename)
            if origin_path is not None:
                orig_img = cv2.imread(str(origin_path), cv2.IMREAD_COLOR)
                if orig_img is not None:
                    o_ih, o_iw = orig_img.shape[:2]
                    orig_ratio = o_iw / o_ih
                    if abs(orig_ratio - box_ratio) >= _RATIO_DIFF_THRESHOLD:
                        source_img = orig_img
                        source_iw, source_ih = o_iw, o_ih
                    else:
                        continue
                else:
                    continue
            else:
                continue

        fi = detect_face_info(None, cv_img=source_img)
        x0, y0, x1, y1 = smart_crop_box(
            source_iw, source_ih, box_w, box_h,
            face_info=fi,
        )

        cropped = source_img[y0:y1, x0:x1]
        success, buf = cv2.imencode(".jpg", cropped, [cv2.IMWRITE_JPEG_QUALITY, 92])
        if not success:
            continue
        new_blob = bytes(buf)

        try:
            blip = shape._element.find(  # type: ignore[union-attr]
                f".//{{{_A}}}blip"
            )
            embed_key = blip.get(  # type: ignore[union-attr]
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
            )
            target_part = shape.part.rels[embed_key].target_part  # type: ignore[attr-defined]
            target_part._blob = new_blob  # type: ignore[attr-defined]
        except Exception:
            continue

        shape.crop_top = 0
        shape.crop_bottom = 0
        shape.crop_left = 0
        shape.crop_right = 0
        changed = True

    return changed
